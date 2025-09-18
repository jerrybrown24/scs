# app.py
# Build: Supply-Chain Simulation Dashboard + 1-click PPTX export
# Sheets expected: Purchasing, Suppliers, Operations, Customers, ROI
# Period field is "rounds" (e.g., -2..6)

import io
import re
import math
import time
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

# --------- Page & constants ---------
st.set_page_config(
    page_title="Supply-Chain Simulation Dashboard",
    page_icon="📦",
    layout="wide"
)

# Fixed thresholds per your spec
THRESH = {
    "reject_gt": 5.0,        # %
    "obsol_gt": 7.0,         # %
    "reliab_lt": 85.0,       # %
    "break_gt": 10.0,        # %
    "unused_gt": 30.0,       # %
    "osa_lt": 92.0,          # %
    "rev_pallet_lt": 40000.0 # currency units
}

# Collect figs for PPT
FIG_REGISTRY = {"ROI": [], "Purchasing": [], "Supply Chain": [], "Operations": [], "Sales": []}

# --------- helpers: parsing & columns ---------
def to_snake(s: str) -> str:
    if not isinstance(s, str):
        s = str(s)
    s = s.strip()
    s = re.sub(r"\s+", " ", s)
    s = s.replace("(", "").replace(")", "").replace("/", " ").replace("%", " pct")
    s = s.lower()
    s = s.replace("  ", " ")
    s = s.replace(" ", "_")
    return s

def normalize_headers(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [to_snake(c) for c in df.columns]
    # preferred synonyms → canonical names
    synonyms = {
        "delivery_reliability": "delivery_reliability_pct",
        "delivery_reliability_pct": "delivery_reliability_pct",
        "rejection": "rejection_pct",
        "obsoletes": "obsolescence_pct",
        "component_bias": "component_bias_pct",
        "stock_pieces": "stock_pieces",
        "stock_value_economic": "stock_value_economic",
        "on_shelf_availability_osa": "osa_pct",
        "osa": "osa_pct",
        "run_time_performance_pct": "run_time_perf_pct",
        "changeover_pct": "changeover_pct",
        "breakdown_pct": "breakdown_pct",
        "unused_capacity_pct": "unused_capacity_pct",
        "overtime_pct": "overtime_pct",
        "start_up_pr": "start_up_production",
        "start_up_production": "start_up_production",
        "gross_margin_per_pallet": "gross_margin_per_pallet",
        "revenue_per_pallet": "revenue_per_pallet",
    }
    ren = {}
    for c in df.columns:
        ren[c] = synonyms.get(c, c)
    df = df.rename(columns=ren)
    # standardize period col
    if "round" in df.columns and "rounds" not in df.columns:
        df = df.rename(columns={"round": "rounds"})
    return df

def find_col(df: pd.DataFrame, *candidates, regex=False):
    cols = list(df.columns)
    if regex:
        pat = re.compile("|".join(candidates))
        for c in cols:
            if pat.search(c):
                return c
        return None
    for c in candidates:
        if c in cols:
            return c
    return None

def parse_percent(series: pd.Series) -> pd.Series:
    """Convert '92.6%' -> 92.6 ; '0.926' -> 92.6 ; numbers stay as % (0-100)."""
    s = series.astype(str).str.strip().str.replace(",", ".", regex=False)
    is_pct = s.str.endswith("%")
    s = np.where(is_pct, s.str[:-1], s)
    out = pd.to_numeric(s, errors="coerce")
    # if values look like fraction (<=1.6) then scale
    if pd.notna(out).sum() and out.dropna().max() <= 1.6:
        out = out * 100.0
    return out

def parse_money(series: pd.Series) -> pd.Series:
    """Remove $, commas, Indian grouping."""
    s = series.astype(str)
    s = s.str.replace("$", "", regex=False)
    s = s.str.replace(",", "", regex=False)
    s = s.str.replace("₹", "", regex=False)
    return pd.to_numeric(s, errors="coerce")

def add_parsed_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    # percent-like columns by name hint
    pct_hints = [c for c in df.columns if c.endswith("_pct") or c.endswith("_percentage") or c.endswith("_percent")]
    for c in pct_hints:
        df[c] = parse_percent(df[c])
    # common explicit names
    for c in ["delivery_reliability_pct","rejection_pct","obsolescence_pct","component_bias_pct",
              "run_time_perf_pct","changeover_pct","breakdown_pct","unused_capacity_pct","overtime_pct",
              "osa_pct"]:
        if c in df.columns:
            df[c] = parse_percent(df[c])
    # currency-like hints
    money_hints = [c for c in df.columns if "revenue" in c or "amount" in c or "value" in c or "margin" in c]
    for c in money_hints:
        if c.endswith("_num"):  # skip already parsed
            continue
        parsed = parse_money(df[c])
        if parsed.notna().any():  # only keep if conversion made sense
            df[c + "_num"] = parsed
    return df

def load_workbook(xfile) -> dict:
    xl = pd.ExcelFile(xfile)
    dfs = {name: add_parsed_columns(normalize_headers(pd.read_excel(xl, sheet_name=name)))
           for name in xl.sheet_names}
    return dfs

def kpi_card(col, label, value, fmt=None, delta=None):
    if fmt is None:
        col.metric(label, value, delta=delta)
    else:
        col.metric(label, fmt.format(value), delta=delta)

def csv_download_button(df, label: str, filename: str):
    st.download_button(
        label=label,
        data=df.to_csv(index=False).encode("utf-8-sig"),
        file_name=filename,
        mime="text/csv"
    )

# --------- PPTX export helpers ---------
def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_section_title(prs: Presentation, title: str):
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    return slide

def add_fig_slide(prs: Presentation, title: str, fig):
    # Use kaleido to convert Plotly fig → PNG bytes
    img = fig.to_image(format="png", scale=2)
    slide = add_section_title(prs, title)
    left = Inches(0.5)
    top = Inches(1.4)
    height = Inches(5.5)
    slide.shapes.add_picture(io.BytesIO(img), left, top, height=height)

def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]  # title & content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

def export_ppt(fig_registry, warning_summary) -> bytes:
    prs = Presentation()
    add_title_slide(prs, "Supply-Chain Dashboard", datetime.now().strftime("%Y-%m-%d %H:%M"))
    # For each tab, dump figs
    for section, items in fig_registry.items():
        if not items:
            continue
        add_section_title(prs, section)
        for title, fig in items:
            try:
                add_fig_slide(prs, f"{section} — {title}", fig)
            except Exception as e:
                add_bullets_slide(prs, f"{section} — {title} (figure)", [f"Could not render figure: {e}"])
    # Warnings & Actions
    bullets = [
        f"Purchasing: {warning_summary.get('Purchasing','0')} red flags (rejection > {THRESH['reject_gt']}% / obsolescence > {THRESH['obsol_gt']}%).",
        f"Supply Chain: {warning_summary.get('Supply Chain','0')} red flags (reliability < {THRESH['reliab_lt']}% / rejection > {THRESH['reject_gt']}%).",
        f"Operations: {warning_summary.get('Operations','0')} red flags (breakdown > {THRESH['break_gt']}% / unused > {THRESH['unused_gt']}%).",
        f"Sales: {warning_summary.get('Sales','0')} red flags (OSA < {THRESH['osa_lt']}% / revenue/pallet < {THRESH['rev_pallet_lt']:.0f}).",
        "Recommended actions:",
        "• Purchasing: expedite QA, review specs, adjust safety stock.",
        "• Supply Chain: renegotiate or dual-source; inspect high-cost lanes.",
        "• Operations: SMED for changeovers; MTBF/RCM to cut breakdowns; re-balance capacity.",
        "• Sales: reset SLAs for low-OSA accounts; revisit margin mix."
    ]
    add_bullets_slide(prs, "Warnings & Actions", bullets)
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

# --------- UI: file upload ---------
st.sidebar.header("📂 Data")
xfile = st.sidebar.file_uploader("Upload Excel workbook (.xlsx)", type=["xlsx", "xlsm"])

with st.sidebar.expander("⚙️ Thresholds (fixed by spec)"):
    st.write(
        f"""
- Rejection > **{THRESH['reject_gt']}%**  
- Obsolescence > **{THRESH['obsol_gt']}%**  
- Supplier reliability < **{THRESH['reliab_lt']}%**  
- Breakdown > **{THRESH['break_gt']}%**, Unused capacity > **{THRESH['unused_gt']}%**  
- OSA < **{THRESH['osa_lt']}%**  
- Revenue / Pallet < **{THRESH['rev_pallet_lt']:.0f}**
"""
    )

if not xfile:
    st.title("📦 Supply-Chain Simulation Dashboard")
    st.info("Upload your Excel workbook (with sheets: Purchasing, Suppliers, Operations, Customers, ROI) to begin.")
    st.stop()

# --------- Load workbook ---------
dfs = load_workbook(xfile)

required = ["purchasing", "suppliers", "operations", "customers", "roi"]
for key in required:
    if key not in [to_snake(k) for k in dfs.keys()]:
        st.warning(f"Sheet named '{key.capitalize()}' not found. Available: {list(dfs.keys())}")

# Map by canonical names
def get_sheet(name):
    for k, v in dfs.items():
        if to_snake(k) == name:
            return v.copy()
    # fallback: first sheet
    return list(dfs.values())[0].copy()

df_p = get_sheet("purchasing")
df_s = get_sheet("suppliers")
df_o = get_sheet("operations")
df_c = get_sheet("customers")
df_r = get_sheet("roi")

# --------- Global rounds filter ---------
def rounds_union():
    cols = []
    for d in [df_p, df_s, df_o, df_c, df_r]:
        if "rounds" in d.columns:
            cols.append(d["rounds"])
    if not cols:
        return []
    return sorted(pd.unique(pd.concat(cols).dropna()))

rounds_all = rounds_union()
st.title("📦 Supply-Chain Simulation Dashboard")
sel_rounds = st.multiselect("Filter by rounds (periods)", rounds_all, default=rounds_all)

# Warning counters for PPT summary
warning_counts = {"Purchasing": 0, "Supply Chain": 0, "Operations": 0, "Sales": 0}

# =========================
# TAB: Overview (ROI)
# =========================
tab_roi, tab_pur, tab_sup, tab_ops, tab_sales = st.tabs(
    ["📈 Overview (ROI)", "🧾 Purchasing", "🚚 Supply Chain", "🏭 Operations", "🛒 Sales & Customers"]
)

with tab_roi:
    FIG_REGISTRY["ROI"].clear()
    st.subheader("Return on Investment & Revenue")

    df = df_r.copy()
    if "rounds" in df.columns and sel_rounds:
        df = df[df["rounds"].isin(sel_rounds)]

    # try to locate ROI & revenue cols
    roi_col = find_col(df, "roi", "roi_pct") or "roi"
    if "roi_pct" not in df.columns and roi_col in df.columns:
        df["roi_pct"] = parse_percent(df[roi_col])

    rev_cols = [c for c in df.columns if c.startswith("realized_revenue")]
    # KPI cards
    c1, c2, c3, c4 = st.columns(4)
    if "roi_pct" in df.columns and not df.empty:
        kpi_card(c1, "Avg ROI %", f"{df['roi_pct'].mean():.2f}%")
        kpi_card(c2, "Min ROI %", f"{df['roi_pct'].min():.2f}%")
        kpi_card(c3, "Max ROI %", f"{df['roi_pct'].max():.2f}%")
    if rev_cols:
        total_rev = df[rev_cols].apply(pd.to_numeric, errors="coerce").sum().sum()
        kpi_card(c4, "Total Realized Revenue", f"{total_rev:,.0f}")

    # 1) ROI line
    if "roi_pct" in df.columns and "rounds" in df.columns and not df.empty:
        fig1 = px.line(df.sort_values("rounds"), x="rounds", y="roi_pct", markers=True, title="ROI % by Round")
        st.plotly_chart(fig1, use_container_width=True)
        FIG_REGISTRY["ROI"].append(("ROI % by Round", fig1))

    # 2) Stacked revenue by round
    if rev_cols and "rounds" in df.columns:
        melt = df[["rounds"] + rev_cols].melt(id_vars="rounds", var_name="revenue_type", value_name="value")
        melt["value"] = pd.to_numeric(melt["value"], errors="coerce")
        fig2 = px.bar(melt, x="rounds", y="value", color="revenue_type", title="Realized Revenue Breakdown", barmode="stack")
        st.plotly_chart(fig2, use_container_width=True)
        FIG_REGISTRY["ROI"].append(("Realized Revenue Breakdown", fig2))

    # 3) Period-over-period change in ROI
    if "roi_pct" in df.columns and "rounds" in df.columns:
        d2 = df.sort_values("rounds")[["rounds","roi_pct"]].copy()
        d2["roi_pop_change"] = d2["roi_pct"].diff()
        fig3 = px.bar(d2, x="rounds", y="roi_pop_change", title="Period-over-Period Change in ROI %")
        st.plotly_chart(fig3, use_container_width=True)
        FIG_REGISTRY["ROI"].append(("PoP Change in ROI %", fig3))

    # 4) ROI histogram
    if "roi_pct" in df.columns:
        fig4 = px.histogram(df, x="roi_pct", nbins=15, title="ROI % Distribution")
        st.plotly_chart(fig4, use_container_width=True)
        FIG_REGISTRY["ROI"].append(("ROI % Distribution", fig4))

    # 5) Waterfall: revenue delta between two rounds
    if rev_cols and "rounds" in df.columns:
        rounds_opts = sorted(df["rounds"].unique().tolist())
        c1, c2 = st.columns(2)
        r_from = c1.selectbox("Revenue delta: from round", rounds_opts, index=0)
        r_to   = c2.selectbox("Revenue delta: to round", rounds_opts, index=len(rounds_opts)-1)
        df_from = df[df["rounds"]==r_from][rev_cols].apply(pd.to_numeric, errors="coerce").sum()
        df_to   = df[df["rounds"]==r_to][rev_cols].apply(pd.to_numeric, errors="coerce").sum()
        delta = (df_to - df_from)
        wf = go.Figure(go.Waterfall(
            x=rev_cols,
            measure=["relative"]*len(rev_cols),
            text=[f"{v:,.0f}" for v in delta],
            y=delta,
        ))
        wf.update_layout(title=f"Revenue Delta by Type — Round {r_from} → {r_to}")
        st.plotly_chart(wf, use_container_width=True)
        FIG_REGISTRY["ROI"].append((f"Revenue Delta {r_from}→{r_to}", wf))

    # 6) Styled table with conditional format (ROI < 0 red)
    if "roi_pct" in df.columns:
        styled = df.copy()
        def color_roi(v):
            color = "background-color: #ffd6d6" if v < 0 else ""
            return [color]*1
        st.dataframe(styled.style.apply(lambda s: ["background-color: #ffd6d6" if x < 0 else "" for x in s] 
                                        if s.name=="roi_pct" else [""]*len(s), axis=0))
    csv_download_button(df, "⬇️ Download ROI (filtered CSV)", "roi_filtered.csv")

# =========================
# TAB: Purchasing (≥7)
# =========================
with tab_pur:
    FIG_REGISTRY["Purchasing"].clear()
    st.subheader("Purchasing KPIs")
    df = df_p.copy()
    if sel_rounds:
        df = df[df.get("rounds").isin(sel_rounds)] if "rounds" in df.columns else df

    comp_col = find_col(df, "component") or "component"
    comps = sorted(df[comp_col].dropna().unique()) if comp_col in df.columns else []
    sel_comps = st.multiselect("Components", comps, default=comps)
    if sel_comps and comp_col in df.columns:
        df = df[df[comp_col].isin(sel_comps)]

    rej = find_col(df, "rejection_pct")
    obs = find_col(df, "obsolescence_pct", "obsoletes")
    deliv = find_col(df, "delivery_reliability_pct", "delivery_reliability")
    purch = find_col(df, "purchase")
    trans = find_col(df, "transport")
    demand = find_col(df, "demand")
    stock_p = find_col(df, "stock_pieces", "stock_pcs")
    stock_val = find_col(df, "stock_value_economic_num", "stock_value_economic")
    comp_bias = find_col(df, "component_bias_pct")
    # 1) Stacked bar: rejection & obsolescence by rounds (facet by component)
    if rej and obs and "rounds" in df.columns and comp_col in df.columns and not df.empty:
        d1 = df[[comp_col,"rounds",rej,obs]].melt(id_vars=[comp_col,"rounds"], var_name="metric", value_name="pct")
        fig = px.bar(d1, x="rounds", y="pct", color="metric", facet_col=comp_col,
                     title="Rejection % & Obsolescence % by Component & Round")
        st.plotly_chart(fig, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Reject & Obsolescence by Component", fig))
    # 2) Delivery reliability trend
    if deliv and comp_col in df.columns and "rounds" in df.columns:
        fig2 = px.line(df.sort_values("rounds"), x="rounds", y=deliv, color=comp_col, markers=True,
                       title="Delivery Reliability % Trend")
        st.plotly_chart(fig2, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Delivery Reliability Trend", fig2))
    # 3) Purchase vs Transport (totals)
    if purch and trans and comp_col in df.columns:
        agg = df.groupby(comp_col)[[purch,trans]].sum().reset_index()
        fig3 = px.bar(agg, x=comp_col, y=[purch,trans], title="Total Purchase vs Transport by Component")
        st.plotly_chart(fig3, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Purchase vs Transport", fig3))
    # 4) Bubble: stock vs economic value (size demand)
    if stock_p and stock_val and demand and comp_col in df.columns:
        fig4 = px.scatter(df, x=stock_p, y=stock_val, size=demand, color=comp_col,
                          title="Inventory Position — Stock vs Economic Value (size=demand)")
        st.plotly_chart(fig4, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Inventory Position", fig4))
    # 5) Component bias %
    if comp_bias and comp_col in df.columns and "rounds" in df.columns:
        fig5 = px.line(df.sort_values("rounds"), x="rounds", y=comp_bias, color=comp_col,
                       title="Component Bias % by Round")
        st.plotly_chart(fig5, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Component Bias %", fig5))
    # 6) Heatmap: reliability vs rejection (mean per component)
    if deliv and rej and comp_col in df.columns:
        heat = df.groupby(comp_col)[[deliv, rej]].mean().round(2)
        fig6 = px.imshow(heat, aspect="auto", title="Heatmap: Mean Reliability vs Rejection by Component")
        st.plotly_chart(fig6, use_container_width=True)
        FIG_REGISTRY["Purchasing"].append(("Heatmap: Reliability vs Rejection", fig6))
    # 7) Warnings table + banners
    warn_df = pd.DataFrame()
    if rej and obs:
        warn_df = df[[comp_col,"rounds",rej,obs]].copy()
        warn_df["warn_reject"] = warn_df[rej] > THRESH["reject_gt"]
        warn_df["warn_obsol"] = warn_df[obs] > THRESH["obsol_gt"]
        st.dataframe(warn_df)
        n_flags = int(warn_df[["warn_reject","warn_obsol"]].sum().sum())
        warning_counts["Purchasing"] = n_flags
        if (warn_df["warn_reject"]).any():
            st.error(f"High Rejection > {THRESH['reject_gt']}% detected.")
        if (warn_df["warn_obsol"]).any():
            st.warning(f"High Obsolescence > {THRESH['obsol_gt']}% detected.")
    csv_download_button(df, "⬇️ Download Purchasing (filtered CSV)", "purchasing_filtered.csv")

# =========================
# TAB: Supply Chain (≥7)
# =========================
with tab_sup:
    FIG_REGISTRY["Supply Chain"].clear()
    st.subheader("Supplier Performance")
    df = df_s.copy()
    if sel_rounds:
        df = df[df.get("rounds").isin(sel_rounds)] if "rounds" in df.columns else df
    sup_col = find_col(df, "supplier") or "supplier"
    sups = sorted(df[sup_col].dropna().unique()) if sup_col in df.columns else []
    sel_sups = st.multiselect("Suppliers", sups, default=sups)
    if sel_sups and sup_col in df.columns:
        df = df[df[sup_col].isin(sel_sups)]

    deliv = find_col(df, "delivery_reliability_pct", "delivery_reliability")
    rej = find_col(df, "rejection_pct", "rejection")
    trans = find_col(df, "transport")
    purch = find_col(df, "purchase")
    orders = find_col(df, "order_lines", "order_line", "order_lines_")
    if not orders:
        # handle "order lines" with space
        orders = find_col(df, "order_lines", regex=True) or find_col(df, "order_lines")
    delivs = find_col(df, "deliveries")

    # 1) Quality bars (facet by supplier)
    if deliv and rej and "rounds" in df.columns and sup_col in df.columns:
        d1 = df[[sup_col,"rounds",deliv,rej]].melt(id_vars=[sup_col,"rounds"], var_name="metric", value_name="pct")
        fig = px.bar(d1, x="rounds", y="pct", color="metric", facet_col=sup_col,
                     title="Reliability vs Rejection by Supplier & Round")
        st.plotly_chart(fig, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Reliability vs Rejection", fig))
    # 2) Order lines trend
    if orders and sup_col in df.columns and "rounds" in df.columns:
        fig2 = px.line(df.sort_values("rounds"), x="rounds", y=orders, color=sup_col, markers=True,
                       title="Order Lines Trend by Supplier")
        st.plotly_chart(fig2, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Order Lines Trend", fig2))
    # 3) Supplier matrix: reliability vs transport
    if deliv and trans and sup_col in df.columns:
        agg = df.groupby(sup_col)[[deliv,trans]].mean().reset_index()
        fig3 = px.scatter(agg, x=deliv, y=trans, text=sup_col, title="Supplier Matrix: Reliability vs Transport Cost")
        fig3.update_traces(textposition="top center")
        st.plotly_chart(fig3, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Reliability vs Transport Matrix", fig3))
    # 4) Purchases vs Deliveries
    if purch and delivs and sup_col in df.columns:
        agg2 = df.groupby(sup_col)[[purch,delivs]].sum().reset_index()
        fig4 = px.bar(agg2, x=sup_col, y=[purch,delivs], title="Total Purchases vs Deliveries by Supplier")
        st.plotly_chart(fig4, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Purchases vs Deliveries", fig4))
    # 5) Funnel: Delivered → Rejected → Accepted
    if delivs and rej and sup_col in df.columns:
        delivered = df.groupby(sup_col)[delivs].sum()
        rej_qty = delivered * (df.groupby(sup_col)[rej].mean() / 100.0)
        accepted = delivered - rej_qty
        f = pd.DataFrame({"supplier": delivered.index,
                          "Delivered": delivered.values,
                          "Rejected": rej_qty.values,
                          "Accepted": accepted.values})
        f_m = f.melt(id_vars="supplier", var_name="stage", value_name="qty")
        fig5 = px.bar(f_m, x="supplier", y="qty", color="stage", title="Flow: Delivered → Rejected → Accepted", barmode="group")
        st.plotly_chart(fig5, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Delivery Flow", fig5))
    # 6) Heatmap: mean KPIs
    metrics = [c for c in [deliv, rej, trans, purch] if c]
    if metrics and sup_col in df.columns:
        heat = df.groupby(sup_col)[metrics].mean().round(2)
        fig6 = px.imshow(heat, aspect="auto", title="Supplier Heatmap: Mean KPIs")
        st.plotly_chart(fig6, use_container_width=True)
        FIG_REGISTRY["Supply Chain"].append(("Heatmap: Supplier KPIs", fig6))
    # 7) Warnings table + banners
    warn_df = pd.DataFrame()
    if deliv and rej:
        warn_df = df[[sup_col,"rounds",deliv,rej]].copy()
        warn_df["warn_low_reliability"] = warn_df[deliv] < THRESH["reliab_lt"]
        warn_df["warn_high_rejection"] = warn_df[rej] > THRESH["reject_gt"]
        st.dataframe(warn_df)
        n_flags = int(warn_df[["warn_low_reliability","warn_high_rejection"]].sum().sum())
        warning_counts["Supply Chain"] = n_flags
        if (warn_df["warn_low_reliability"]).any():
            st.error(f"Low delivery reliability < {THRESH['reliab_lt']}% detected.")
        if (warn_df["warn_high_rejection"]).any():
            st.warning(f"High supplier rejection > {THRESH['reject_gt']}% detected.")
    csv_download_button(df, "⬇️ Download Suppliers (filtered CSV)", "suppliers_filtered.csv")

# =========================
# TAB: Operations (≥7)
# =========================
with tab_ops:
    FIG_REGISTRY["Operations"].clear()
    st.subheader("Operations — Bottling Lines")
    df = df_o.copy()
    if sel_rounds:
        df = df[df.get("rounds").isin(sel_rounds)] if "rounds" in df.columns else df
    line_col = find_col(df, "bottling_line", "bottling", "line") or "bottling_line"
    lines = sorted(df[line_col].dropna().unique()) if line_col in df.columns else []
    sel_lines = st.multiselect("Production Lines", lines, default=lines)
    if sel_lines and line_col in df.columns:
        df = df[df[line_col].isin(sel_lines)]

    runp = find_col(df, "run_time_perf_pct")
    chg = find_col(df, "changeover_pct")
    brk = find_col(df, "breakdown_pct")
    unused = find_col(df, "unused_capacity_pct")
    ot = find_col(df, "overtime_pct")
    start = find_col(df, "start_up_production")
    prod = find_col(df, "production")

    # 1) Runtime performance trend
    if runp and "rounds" in df.columns:
        fig1 = px.line(df.sort_values("rounds"), x="rounds", y=runp, color=line_col, markers=True,
                       title="Run Time Performance % by Round")
        st.plotly_chart(fig1, use_container_width=True)
        FIG_REGISTRY["Operations"].append(("Run Time Performance %", fig1))
    # 2) Losses stacked
    parts = [x for x in [chg, brk, unused] if x]
    if parts and "rounds" in df.columns and line_col in df.columns:
        melt = df[[line_col,"rounds"]+parts].melt(id_vars=[line_col,"rounds"], var_name="loss", value_name="pct")
        fig2 = px.bar(melt, x="rounds", y="pct", color="loss", facet_col=line_col, barmode="stack",
                      title="Losses — Changeover / Breakdown / Unused Capacity")
        st.plotly_chart(fig2, use_container_width=True)
        FIG_REGISTRY["Operations"].append(("Losses (Stacked)", fig2))
    # 3) Donut: Overtime vs Normal
    if ot and runp:
        df["normal_run_pct"] = np.clip(df[runp] - df[ot], 0, None)
        donut = df[["normal_run_pct", ot]].sum()
        donut_df = pd.DataFrame({"type": ["Normal Run", "Overtime"], "pct": [donut["normal_run_pct"], donut[ot]]})
        fig3 = px.pie(donut_df, names="type", values="pct", hole=0.55, title="Overtime vs Normal Run (Aggregated)")
        st.plotly_chart(fig3, use_container_width=True)
        FIG_REGISTRY["Operations"].append(("Overtime vs Normal Run", fig3))
    # 4) Production vs Start-up
    if prod and start and line_col in df.columns:
        agg = df.groupby(line_col)[[prod, start]].sum().reset_index()
        fig4 = px.bar(agg, x=line_col, y=[prod, start], title="Production Output vs Start-up Production")
        st.plotly_chart(fig4, use_container_width=True)
        FIG_REGISTRY["Operations"].append(("Production vs Start-up", fig4))
    # 5) Heatmap: efficiency vs breakdown by round
    if runp and brk and "rounds" in df.columns:
        h = df.groupby("rounds")[[runp, brk]].mean().round(2)
        fig5 = px.imshow(h.T, aspect="auto", title="Heatmap: Efficiency vs Breakdown by Round")
        st.plotly_chart(fig5, use_container_width=True)
        FIG_REGISTRY["Operations"].append(("Heatmap: Efficiency vs Breakdown", fig5))
    # 6) KPI cards
    c1, c2, c3 = st.columns(3)
    if runp: kpi_card(c1, "Avg Run Time Perf %", f"{df[runp].mean():.1f}%")
    if brk:  kpi_card(c2, "Avg Breakdown %", f"{df[brk].mean():.1f}%")
    if unused: kpi_card(c3, "Avg Unused Capacity %", f"{df[unused].mean():.1f}%")
    # 7) Warnings table + banners
    warn_df = pd.DataFrame()
    if brk and unused:
        warn_df = df[[line_col,"rounds",brk,unused]].copy()
        warn_df["warn_breakdown"] = warn_df[brk] > THRESH["break_gt"]
        warn_df["warn_unused"] = warn_df[unused] > THRESH["unused_gt"]
        st.dataframe(warn_df)
        n_flags = int(warn_df[["warn_breakdown","warn_unused"]].sum().sum())
        warning_counts["Operations"] = n_flags
        if (warn_df["warn_breakdown"]).any():
            st.error(f"High breakdown > {THRESH['break_gt']}% detected.")
        if (warn_df["warn_unused"]).any():
            st.warning(f"High unused capacity > {THRESH['unused_gt']}% detected.")
    csv_download_button(df, "⬇️ Download Operations (filtered CSV)", "operations_filtered.csv")

# =========================
# TAB: Sales & Customers (≥7)
# =========================
with tab_sales:
    FIG_REGISTRY["Sales"].clear()
    st.subheader("Sales & Customers")
    df = df_c.copy()
    if sel_rounds:
        df = df[df.get("rounds").isin(sel_rounds)] if "rounds" in df.columns else df

    cust = find_col(df, "customer") or "customer"
    customers = sorted(df[cust].dropna().unique()) if cust in df.columns else []
    sel_customers = st.multiselect("Customers", customers, default=customers)
    if sel_customers and cust in df.columns:
        df = df[df[cust].isin(sel_customers)]

    # detect service levels
    svc_cols = [c for c in df.columns if c.startswith("service_lev")]
    revp = find_col(df, "revenue_per_pallet_num", "revenue_per_pallet")
    gmp = find_col(df, "gross_margin_per_pallet_num", "gross_margin_per_pallet")
    shipments = find_col(df, "shipments")
    orderlines = find_col(df, "order_lines", "order_line")
    if not orderlines:
        orderlines = find_col(df, "order", regex=True)
    osa = find_col(df, "osa_pct", "osa")
    attained = find_col(df, "attained_s", "attained_service")

    # 1) Service levels stacked/grouped
    if svc_cols and cust in df.columns and "rounds" in df.columns:
        m = df[[cust, "rounds"] + svc_cols].melt(id_vars=[cust,"rounds"], var_name="service_tier", value_name="pct")
        fig1 = px.bar(m, x="rounds", y="pct", color="service_tier", facet_col=cust, barmode="group",
                      title="Service Levels by Customer")
        st.plotly_chart(fig1, use_container_width=True)
        FIG_REGISTRY["Sales"].append(("Service Levels", fig1))
    # 2) Bubble: rev vs margin (size shipments)
    if revp and gmp and shipments and cust in df.columns:
        fig2 = px.scatter(df, x=revp, y=gmp, size=shipments, color=cust,
                          title="Revenue per Pallet vs Gross Margin per Pallet")
        st.plotly_chart(fig2, use_container_width=True)
        FIG_REGISTRY["Sales"].append(("Revenue vs Margin Bubble", fig2))
    # 3) Orders vs Shipments
    if shipments and orderlines and cust in df.columns and "rounds" in df.columns:
        agg = df.groupby([cust,"rounds"])[[shipments,orderlines]].sum().reset_index()
        fig3 = px.bar(agg, x="rounds", y=[shipments,orderlines], color=cust, barmode="group",
                      title="Orders vs Shipments (by Round & Customer)")
        st.plotly_chart(fig3, use_container_width=True)
        FIG_REGISTRY["Sales"].append(("Orders vs Shipments", fig3))
    # 4) OSA trend
    if osa and cust in df.columns and "rounds" in df.columns:
        fig4 = px.line(df.sort_values("rounds"), x="rounds", y=osa, color=cust, markers=True,
                       title="On-Shelf Availability (OSA) % Trend")
        st.plotly_chart(fig4, use_container_width=True)
        FIG_REGISTRY["Sales"].append(("OSA % Trend", fig4))
    # 5) Heatmap: order lines vs attained service
    if orderlines and attained and cust in df.columns:
        heat = df.groupby(cust)[[orderlines, attained]].mean().round(2)
        fig5 = px.imshow(heat, aspect="auto", title="Heatmap: Order Lines vs Attained Service")
        st.plotly_chart(fig5, use_container_width=True)
        FIG_REGISTRY["Sales"].append(("Heatmap: Orders vs Attained Service", fig5))
    # 6) KPI cards
    c1, c2, c3 = st.columns(3)
    if revp: kpi_card(c1, "Avg Revenue / Pallet", f"{pd.to_numeric(df[revp], errors='coerce').mean():,.0f}")
    if gmp:  kpi_card(c2, "Avg Gross Margin / Pallet", f"{pd.to_numeric(df[gmp], errors='coerce').mean():,.0f}")
    if osa:  kpi_card(c3, "Avg OSA %", f"{df[osa].mean():.1f}%")
    # 7) Warnings table + banners
    warn_df = pd.DataFrame()
    if osa and revp:
        rnum = pd.to_numeric(df[revp], errors="coerce")
        warn_df = df[[cust,"rounds",osa]].copy()
        warn_df["revenue_per_pallet_num"] = rnum
        warn_df["warn_low_osa"] = warn_df[osa] < THRESH["osa_lt"]
        warn_df["warn_low_rev"] = warn_df["revenue_per_pallet_num"] < THRESH["rev_pallet_lt"]
        st.dataframe(warn_df)
        n_flags = int(warn_df[["warn_low_osa","warn_low_rev"]].sum().sum())
        warning_counts["Sales"] = n_flags
        if (warn_df["warn_low_osa"]).any():
            st.error(f"Low OSA < {THRESH['osa_lt']}% detected.")
        if (warn_df["warn_low_rev"]).any():
            st.warning(f"Low Revenue/Pallet < {THRESH['rev_pallet_lt']:.0f} detected.")
    csv_download_button(df, "⬇️ Download Customers (filtered CSV)", "customers_filtered.csv")

# --------- PPT export (sidebar) ---------
st.sidebar.header("📤 Export")
if st.sidebar.button("Export current view to PowerPoint (.pptx)"):
    try:
        deck_bytes = export_ppt(FIG_REGISTRY, warning_counts)
        st.sidebar.success("Deck generated.")
        st.sidebar.download_button(
            "Download SupplyChain_Dashboard.pptx",
            data=deck_bytes,
            file_name=f"SupplyChain_Dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.sidebar.error(f"Export failed: {e}")
